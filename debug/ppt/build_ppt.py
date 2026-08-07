"""Genera la PPT de avance de DetoVision (Enaex) a partir de las imagenes de debug/ppt.

Reutiliza el estilo del reporte de batimetria (reporte_7-B_local.pptx): tema,
layouts y logo KAUEL. Cada tronadura ocupa 2 laminas:

  A) mascara de cambios | mascara con trayectorias detectadas  (+ KPIs del HUD)
  B) placeholder de PowerPoint (imagen o video, marco 16:9) + observaciones

Uso:
    uv run --no-project --with python-pptx --with pillow python debug/ppt/build_ppt.py
"""

from __future__ import annotations

import shutil
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
TEMPLATE = HERE / "reporte_7-B_local.pptx"
ASSETS = HERE / "_assets"
OUTPUT = HERE / "DetoVision_Avance_Enaex_2026-07-24.pptx"

# ---------------------------------------------------------------- estilo
NAVY = RGBColor(0x14, 0x2C, 0x44)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xEE, 0xF1, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1B, 0x24, 0x2E)
FONT = "Arial"

SW, SH = 13.3333, 7.5  # pulgadas
BAND_H = 1.12

FECHA_PRESENTACION = "Viernes 24 de julio de 2026"
FECHA_GENERACION = "22-07-2026"

# --------------------------------------------------- datos por tronadura
# (id, frame de referencia, umbral de velocidad, rocas detectadas) -- leidos
# del HUD de cada mascara_voladura_analisis_*.jpg
TRONADURAS = [
    {"id": "03", "frame": "843", "vmin": "3.6", "rocks": "570"},
    {"id": "04", "frame": "774", "vmin": "10.2", "rocks": "3.736"},
    {"id": "07", "frame": "827", "vmin": "5.0", "rocks": "1.153"},
    {"id": "09", "frame": "947", "vmin": "1.7", "rocks": "2.041"},
    {"id": "11", "frame": "787", "vmin": "4.2", "rocks": "1.073"},
    {"id": "13", "frame": "1050", "vmin": "0.9", "rocks": "19.482"},
]

PLACEHOLDER_TXT = "[completar]"
PLACEHOLDER_COLOR = RGBColor(0xB0, 0x86, 0x2A)  # ámbar apagado: se ve, no grita


# ------------------------------------------------------------- utilidades
def preparar_assets() -> None:
    """Extrae el logo de la plantilla y reescala las imagenes 4K a un peso sano."""
    ASSETS.mkdir(exist_ok=True)

    logo = ASSETS / "logo_kauel.png"
    if not logo.exists():
        with zipfile.ZipFile(TEMPLATE) as z:
            logo.write_bytes(z.read("ppt/media/image2.png"))

    def reescalar(src: Path, dst: Path, ancho: int = 2400) -> None:
        if dst.exists() or not src.exists():
            return
        im = Image.open(src)
        modo = "L" if im.mode == "L" else "RGB"
        im = im.convert(modo)
        if im.width > ancho:
            im = im.resize((ancho, round(im.height * ancho / im.width)), Image.LANCZOS)
        im.save(dst, quality=90, subsampling=0)

    for t in TRONADURAS:
        n = int(t["id"])
        reescalar(HERE / f"mascara_cambios_final_sinbin_{n}.png", ASSETS / f"mask_{n}.jpg")
        reescalar(HERE / f"mascara_voladura_analisis_{n}.jpg", ASSETS / f"tray_{n}.jpg")

    align = HERE.parent / "out" / "3_secuencia_filtro" / "alineacion_tiros.png"
    reescalar(align, ASSETS / "alineacion_tiros.jpg", 1800)


def preparar_plantilla() -> Presentation:
    """Abre la plantilla, borra sus laminas y habilita un placeholder de contenido."""
    prs = Presentation(TEMPLATE)

    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.get(qn("r:id")))
        sld_id_lst.remove(sld_id)

    # El layout "Picture with Caption" trae dos placeholders 'body' sin indice.
    # Lo convertimos en: titulo + UN placeholder de contenido (idx=1). Un
    # placeholder de contenido vacio es el que muestra en PowerPoint los iconos
    # de "insertar imagen / video", y ajusta el archivo al marco.
    layout = prs.slide_layouts[10]
    bodies = []
    for shape in list(layout.shapes):
        phs = shape.element.findall(".//" + qn("p:ph"))
        if not phs:
            continue
        ph = phs[0]
        if ph.get("type") == "body":
            bodies.append((shape, ph))
    for i, (shape, ph) in enumerate(bodies):
        if i == 0:
            if "type" in ph.attrib:
                del ph.attrib["type"]
            ph.set("idx", "1")
        else:
            shape.element.getparent().remove(shape.element)
    return prs


def add_slide(prs: Presentation, con_placeholder: bool = False):
    """Lamina en blanco (o con placeholder de contenido), sin el placeholder de titulo."""
    layout = prs.slide_layouts[10 if con_placeholder else 8]
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.placeholders):
        if shape.placeholder_format.type is not None and "TITLE" in str(
            shape.placeholder_format.type
        ):
            shape.element.getparent().remove(shape.element)
        elif not con_placeholder:
            shape.element.getparent().remove(shape.element)
    return slide


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.text_frame.text = ""
    return shp


def txt(
    slide,
    x,
    y,
    w,
    h,
    texto,
    size=12,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    space_after=6,
    line_spacing=None,
):
    """Caja de texto. `texto` puede ser str o lista de lineas."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lineas = [texto] if isinstance(texto, str) else list(texto)
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = linea
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def banda(slide, titulo, derecha=""):
    rect(slide, 0, 0, SW, BAND_H, fill=NAVY)
    rect(slide, 0, BAND_H, SW, 0.06, fill=TEAL)
    txt(slide, 0.55, 0.16, 9.4, 0.85, titulo, size=24, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if derecha:
        txt(
            slide,
            9.9,
            0.16,
            2.9,
            0.85,
            derecha,
            size=11,
            color=RGBColor(0xB8, 0xC4, 0xD0),
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def pie(slide, numero=None):
    logo = ASSETS / "logo_kauel.png"
    slide.shapes.add_picture(str(logo), Inches(11.55), Inches(6.92), width=Inches(1.35))
    if numero is not None:
        txt(slide, 0.55, 6.98, 1.5, 0.3, str(numero), size=10, color=GREY)


def imagen(slide, ruta: Path, x, y, w):
    """Inserta una imagen con ancho fijo y devuelve el shape."""
    return slide.shapes.add_picture(str(ruta), Inches(x), Inches(y), width=Inches(w))


def bullets(slide, x, y, w, h, items, size=14, color=DARK, gap=10):
    """Viñetas nativas con sangría francesa. Las líneas que empiezan con el
    marcador de placeholder se pintan en otro color, para ubicarlas de un vistazo."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    sangria = Inches(0.26)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.2
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(sangria))
        pPr.set("indent", str(-sangria))
        for tag, attrs in ((qn("a:buFont"), {"typeface": FONT}), (qn("a:buChar"), {"char": "•"})):
            pPr.append(pPr.makeelement(tag, attrs))
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = PLACEHOLDER_COLOR if item.startswith(PLACEHOLDER_TXT) else color
    return box


# --------------------------------------------------------------- laminas
def lamina_portada(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, 0.32, SH, fill=TEAL)
    txt(slide, 0.95, 2.05, 11.5, 0.7, "DETOVISION", size=20, bold=True, color=TEAL)
    txt(
        slide,
        0.9,
        2.55,
        11.5,
        1.7,
        "Detección de flyrocks\nen tronaduras",
        size=40,
        bold=True,
        color=NAVY,
        space_after=0,
        line_spacing=1.05,
    )
    rect(slide, 0.95, 4.42, 2.2, 0.045, fill=TEAL)
    txt(slide, 0.95, 4.72, 11.0, 0.4, "Enaex  ·  Avance de entrega", size=17, color=DARK)
    txt(slide, 0.95, 5.22, 11.0, 0.35, FECHA_PRESENTACION, size=13, color=GREY)
    txt(slide, 0.95, 6.62, 6.0, 0.3, f"Generado {FECHA_GENERACION}", size=10, color=GREY)
    pie(slide)
    return slide


def lamina_contexto(prs, n):
    slide = add_slide(prs)
    banda(slide, "Contexto y objetivo")
    bullets(
        slide,
        0.55,
        1.75,
        7.6,
        4.4,
        [
            "Objetivo: detectar y visualizar las trayectorias de flyrocks a partir "
            "del video de dron de cada tronadura.",
            "Insumo: video aéreo 4K, una tronadura por video.",
            "Esta entrega: máscara visual de trayectorias por tronadura. Las métricas "
            "por roca (alcance en metros, zonas) vienen en la etapa siguiente.",
            "Alcance del avance: 6 tronaduras procesadas.",
            f"{PLACEHOLDER_TXT} Agregar contexto de faena / alcance del piloto.",
        ],
        size=15,
    )
    rect(slide, 8.5, 1.75, 4.3, 3.05, fill=LIGHT)
    txt(slide, 8.85, 2.05, 3.6, 0.35, "EN ESTA PRESENTACIÓN", size=11, bold=True, color=TEAL)
    txt(
        slide,
        8.85,
        2.5,
        3.6,
        2.2,
        [
            "1.   Resultados por tronadura",
            "2.   Resumen de la entrega",
            "3.   Estado actual y brechas",
            "4.   Próximos pasos",
        ],
        size=13,
        color=NAVY,
        space_after=10,
    )
    pie(slide, n)
    return slide


def lamina_como_leer(prs, n):
    slide = add_slide(prs)
    banda(slide, "Cómo leer estas láminas")
    w = 6.15
    x1, x2 = 0.35, 6.83
    y = 2.0
    txt(slide, x1, 1.58, w, 0.35, "1 · MÁSCARA DE CAMBIOS", size=12, bold=True, color=NAVY)
    txt(slide, x2, 1.58, w, 0.35, "2 · TRAYECTORIAS DETECTADAS", size=12, bold=True, color=NAVY)
    imagen(slide, ASSETS / "mask_7.jpg", x1, y, w)
    imagen(slide, ASSETS / "tray_7.jpg", x2, y, w)
    cap_y = y + w * 9 / 16 + 0.18
    txt(
        slide,
        x1,
        cap_y,
        w,
        1.0,
        "Acumula todo el movimiento registrado durante el evento. Las estelas claras "
        "que salen del centro son las rocas proyectadas.",
        size=12,
        color=GREY,
        line_spacing=1.2,
    )
    txt(
        slide,
        x2,
        cap_y,
        w,
        1.0,
        "Cada color es una trayectoria reconstruida por el algoritmo: un mismo "
        "fragmento seguido a lo largo de los cuadros del video.",
        size=12,
        color=GREY,
        line_spacing=1.2,
    )
    pie(slide, n)
    return slide


def kpi(slide, x, y, w, h, valor, etiqueta):
    rect(slide, x, y, w, h, fill=LIGHT)
    rect(slide, x, y, 0.06, h, fill=TEAL)
    txt(slide, x + 0.3, y + 0.14, w - 0.5, 0.45, valor, size=20, bold=True, color=NAVY)
    txt(slide, x + 0.3, y + 0.6, w - 0.5, 0.3, etiqueta, size=10, color=GREY)


def lamina_tronadura_a(prs, t, n):
    slide = add_slide(prs)
    banda(slide, f"Tronadura {t['id']}", "Vista general")
    w = 6.15
    x1, x2 = 0.35, 6.83
    y = 1.85
    txt(slide, x1, 1.45, w, 0.35, "1 · MÁSCARA DE CAMBIOS", size=11, bold=True, color=NAVY)
    txt(slide, x2, 1.45, w, 0.35, "2 · TRAYECTORIAS DETECTADAS", size=11, bold=True, color=NAVY)
    imagen(slide, ASSETS / f"mask_{int(t['id'])}.jpg", x1, y, w)
    imagen(slide, ASSETS / f"tray_{int(t['id'])}.jpg", x2, y, w)

    ky, kh, kw = 5.72, 0.95, 4.0
    kpi(slide, 0.35, ky, kw, kh, t["frame"], "CUADRO DE REFERENCIA DEL VIDEO")
    kpi(slide, 4.65, ky, kw, kh, f"≥ {t['vmin']}", "UMBRAL DE VELOCIDAD (px/cuadro)")
    kpi(slide, 8.95, ky, kw, kh, t["rocks"], "TRAYECTORIAS DETECTADAS")
    pie(slide, n)
    return slide


def lamina_tronadura_b(prs, t, n):
    slide = add_slide(prs, con_placeholder=True)
    banda(slide, f"Tronadura {t['id']} — Zona de interés", "Detalle ampliado")

    px, py, pw = 0.35, 1.62, 7.9
    ph_h = pw * 9 / 16  # marco 16:9 exacto: no hay que recortar nada

    # Guía punteada AL FONDO: hace visible el marco mientras está vacío y queda
    # tapada apenas se inserta la imagen o el video.
    guia = rect(slide, px, py, pw, ph_h, fill=RGBColor(0xF7, 0xF9, 0xFB), line=RGBColor(0xC3, 0xCB, 0xD4))
    guia.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    spTree = slide.shapes._spTree
    spTree.remove(guia._element)
    spTree.insert(2, guia._element)

    marco = next((s for s in slide.placeholders), None)
    if marco is not None:
        marco.left, marco.top = Inches(px), Inches(py)
        marco.width, marco.height = Inches(pw), Inches(ph_h)

    txt(
        slide,
        px,
        py + ph_h + 0.12,
        pw,
        0.6,
        "Clic en el ícono del recuadro → Imagen o Video. El marco ya está en 16:9: "
        "el archivo entra sin recortar ni reajustar.",
        size=10,
        color=GREY,
        line_spacing=1.2,
    )

    ox, ow = 8.55, 4.28
    rect(slide, ox, py, ow, 4.44, fill=LIGHT)
    rect(slide, ox, py, ow, 0.42, fill=NAVY)
    txt(slide, ox + 0.28, py + 0.05, ow - 0.5, 0.32, "OBSERVACIONES", size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    bullets(
        slide,
        ox + 0.28,
        py + 0.72,
        ow - 0.56,
        3.5,
        [
            f"{PLACEHOLDER_TXT} Qué zona se está ampliando.",
            f"{PLACEHOLDER_TXT} Qué se observa (alcance, dirección, concentración).",
            f"{PLACEHOLDER_TXT} Implicancia operacional o acción.",
        ],
        size=12,
        color=GREY,
        gap=14,
    )
    pie(slide, n)
    return slide


def lamina_tabla(prs, n):
    slide = add_slide(prs)
    banda(slide, "Resumen de tronaduras procesadas")
    filas, cols = len(TRONADURAS) + 1, 4
    tabla = slide.shapes.add_table(
        filas, cols, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.55 * filas)
    ).table
    tabla.columns[0].width = Inches(2.6)
    for c in range(1, cols):
        tabla.columns[c].width = Inches(2.96)

    encabezados = ["TRONADURA", "CUADRO DE REFERENCIA", "UMBRAL DE VELOCIDAD", "TRAYECTORIAS DETECTADAS"]
    datos = [encabezados] + [
        [f"Tronadura {t['id']}", t["frame"], f"≥ {t['vmin']}", t["rocks"]] for t in TRONADURAS
    ]
    for r, fila in enumerate(datos):
        tabla.rows[r].height = Inches(0.52)
        for c, valor in enumerate(fila):
            celda = tabla.cell(r, c)
            celda.text = valor
            celda.vertical_anchor = MSO_ANCHOR.MIDDLE
            celda.margin_left = Inches(0.18)
            celda.fill.solid()
            celda.fill.fore_color.rgb = NAVY if r == 0 else (WHITE if r % 2 else LIGHT)
            p = celda.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.name = FONT
            run.font.size = Pt(10 if r == 0 else 13)
            run.font.bold = r == 0 or c == 0
            run.font.color.rgb = WHITE if r == 0 else NAVY

    txt(
        slide,
        0.9,
        1.7 + 0.55 * filas + 0.35,
        11.5,
        0.8,
        f"{PLACEHOLDER_TXT} Nota: el umbral de velocidad se ajusta por video. Valores muy bajos "
        "dejan entrar ruido — ver «Estado actual y brechas».",
        size=11,
        color=GREY,
    )
    pie(slide, n)
    return slide


def lamina_entrega(prs, n):
    slide = add_slide(prs)
    banda(slide, "Qué se entrega el viernes")
    bullets(
        slide,
        0.55,
        1.85,
        12.2,
        4.4,
        [
            "Máscara de cambios y máscara con trayectorias detectadas, para las 6 tronaduras.",
            f"{PLACEHOLDER_TXT} Video con las trayectorias dibujadas sobre el material original.",
            f"{PLACEHOLDER_TXT} Formato y vía de entrega (carpeta, informe, plataforma).",
            f"{PLACEHOLDER_TXT} Agregar / profundizar.",
        ],
        size=16,
        gap=16,
    )
    pie(slide, n)
    return slide


def lamina_brechas(prs, n):
    slide = add_slide(prs)
    banda(slide, "Estado actual y brechas")
    txt(slide, 0.55, 1.6, 7.6, 0.35, "LO QUE YA FUNCIONA", size=11, bold=True, color=TEAL)
    bullets(
        slide,
        0.55,
        2.0,
        7.6,
        1.7,
        [
            "Compensación del movimiento del dron y detección del movimiento real.",
            "Reconstrucción de trayectorias con continuidad física (sin saltos imposibles).",
        ],
        size=13,
        color=DARK,
        gap=8,
    )
    txt(slide, 0.55, 3.75, 7.6, 0.35, "LO QUE FALTA", size=11, bold=True, color=RGBColor(0xC0, 0x50, 0x4D))
    bullets(
        slide,
        0.55,
        4.15,
        7.6,
        2.2,
        [
            "Sobre-detección en algunos videos (Tronadura 13: 19.482 trazas; Tronadura 04: "
            "3.736): con umbral bajo entra ruido de fondo y humo.",
            "Separar roca de humo en la zona densa central, donde las trayectorias se cruzan.",
            "Calibración métrica: hoy todo está en píxeles, aún no en metros.",
            f"{PLACEHOLDER_TXT} Agregar.",
        ],
        size=13,
        color=DARK,
        gap=8,
    )
    rect(slide, 8.5, 1.6, 4.3, 4.75, fill=LIGHT)
    txt(slide, 8.85, 1.95, 3.6, 1.0, "El cuello de botella no es\nver el movimiento: es\ndecidir qué es roca.", size=16, bold=True, color=NAVY, space_after=0, line_spacing=1.25)
    txt(
        slide,
        8.85,
        3.3,
        3.6,
        2.7,
        "La señal está en el video — las estelas se ven. Lo que falta es el criterio "
        "que distingue una roca de una voluta de humo. Ahí es donde el conocimiento "
        "del operador entra al algoritmo.",
        size=12,
        color=GREY,
        line_spacing=1.3,
    )
    pie(slide, n)
    return slide


def lamina_proximos(prs, n):
    slide = add_slide(prs)
    banda(slide, "Próximos pasos")
    bullets(
        slide,
        0.55,
        1.75,
        6.5,
        4.4,
        [
            "Incorporar la secuencia de tiros: usar el orden y tiempo de detonación para "
            "descartar lo que no puede venir de un tiro ya detonado.",
            "Calibración en terreno: ubicar la malla de tiros sobre la imagen para pasar de "
            "píxeles a metros y medir alcances reales.",
            "Revisión asistida: interfaz para validar, unir o descartar trayectorias antes "
            "del informe final.",
            "Reporte por tronadura: conteo, alcance máximo y zonas alcanzadas.",
            f"{PLACEHOLDER_TXT} Agregar / priorizar con el cliente.",
        ],
        size=14,
        gap=12,
    )
    txt(slide, 7.3, 1.75, 5.5, 0.35, "EJEMPLO — MALLA DE TIROS SOBRE EL EVENTO", size=10, bold=True, color=NAVY)
    imagen(slide, ASSETS / "alineacion_tiros.jpg", 7.3, 2.15, 5.5)
    txt(
        slide,
        7.3,
        2.15 + 5.5 * 9 / 16 + 0.15,
        5.5,
        0.8,
        "Cada punto es un tiro, coloreado por su tiempo de detonación (azul = primero, "
        "rojo = último). Sirve para exigir que cada roca salga de un tiro ya detonado.",
        size=10,
        color=GREY,
        line_spacing=1.2,
    )
    pie(slide, n)
    return slide


def lamina_cierre(prs, n):
    slide = add_slide(prs)
    rect(slide, 0, 0, SW, SH, fill=NAVY)
    rect(slide, 0, 0, 0.32, SH, fill=TEAL)
    txt(slide, 1.2, 3.05, 10.5, 0.9, "Preguntas y comentarios", size=32, bold=True, color=WHITE)
    rect(slide, 1.25, 4.15, 2.2, 0.045, fill=TEAL)
    txt(slide, 1.25, 4.45, 10.5, 0.4, "DetoVision  ·  Detección de flyrocks", size=14, color=RGBColor(0xB8, 0xC4, 0xD0))
    slide.shapes.add_picture(str(ASSETS / "logo_kauel.png"), Inches(11.4), Inches(6.75), width=Inches(1.45))
    return slide


# ------------------------------------------------------------------ main
def main() -> None:
    preparar_assets()
    prs = preparar_plantilla()

    n = 1
    lamina_portada(prs)
    for constructor in (lamina_contexto, lamina_como_leer):
        n += 1
        constructor(prs, n)
    for t in TRONADURAS:
        n += 1
        lamina_tronadura_a(prs, t, n)
        n += 1
        lamina_tronadura_b(prs, t, n)
    for constructor in (lamina_tabla, lamina_entrega, lamina_brechas, lamina_proximos):
        n += 1
        constructor(prs, n)
    n += 1
    lamina_cierre(prs, n)

    prs.save(OUTPUT)
    mb = OUTPUT.stat().st_size / 1e6
    print(f"OK  {OUTPUT.name}  —  {len(prs.slides.__iter__.__self__._sldIdLst)} láminas, {mb:.1f} MB")


if __name__ == "__main__":
    main()
